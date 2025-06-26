from pptx import Presentation
import os

def read_pptx(file_path):
    """
    讀取 PPTX 檔案並提取所有幻燈片的文字內容
    :param file_path: PPTX 檔案路徑
    :return: 幻燈片內容的列表
    """
    prs = Presentation(file_path)
    slides_content = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_content.append(f"Slide {slide_number}:\n" + "\n".join(slide_text))

    return slides_content

def extract_images_from_pptx(pptx_path):
    """
    從 PPTX 檔案中提取所有圖片並儲存
    :param pptx_path: PPTX 檔案路徑
    :param output_folder: 圖片儲存資料夾
    """
    prs = Presentation(pptx_path)

    if not os.path.exists("extracted_img"):
        os.makedirs("extracted_img")
        output_folder = "extracted_img"
    else:
        output_folder = "extracted_img"

    image_count = 0
    for slide_number, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if hasattr(shape, "image"):
                image = shape.image
                image_bytes = image.blob
                image_format = image.ext  # 取得圖片格式 (如 jpg, png, etc.)

                image_filename = os.path.join(output_folder, f"slide_page_{slide_number}_{image_count}.{image_format}")
                with open(image_filename, "wb") as img_file:
                    img_file.write(image_bytes)

                print(f"已儲存圖片: {image_filename}")
                image_count += 1

if __name__ == "__main__":
    pptx_file = "sample\sample_ppt.pptx"  # 請替換為你的 PPTX 檔案路徑
    content = read_pptx(pptx_file)
    extract_images_from_pptx(pptx_file)

    for slide_text in content:
        print(slide_text)
        print("=" * 50)
        with open("extracted_txt/ppt_output.txt", "a+", encoding="utf-8") as txt_file:
            txt_file.write(slide_text + "\n" + "=" * 50 + "\n")
    